from pptx import Presentation  
from pptx.oxml.ns import qn  
from lxml import etree  
  
prs = Presentation('DeepSeek能力全景.pptx')  
for i, slide in enumerate(prs.slides):  
    print(f'\n=== 第{i+1}页 ===')  
    for shape in slide.shapes:  
        cNvPr = shape._element.find(qn('p:cNvPr'))  
        sid = cNvPr.get('id') if cNvPr is not None else 'N/A'  
        print(f'  [{sid}] {shape.shape_type} - {shape.name}')  
