"""
📊 PPT制作技能 v4.0 — 双引擎：python-pptx(优先) + 纯Python回退
"""
import os, json
from datetime import datetime
from dataclasses import dataclass
from typing import Dict, Any, List
from pyclaw.skill import SkillMetadata
from pyclaw.pyclaw_types import ToolDefinition, ToolResult

# ── 引擎检测 ──
_HAS_PPTX = False
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _HAS_PPTX = True
except ImportError:
    pass

SKILL_CLASS = "PPTSkill"

LAYOUTS = {
    "title":         "深色封面",
    "title_center":  "深色居中封面",
    "content_light": "浅色内容",
    "content_dark":  "深色内容",
    "two_column":    "浅色双栏",
    "two_column_dark":"深色双栏",
    "features_dark": "深色卡片",
    "process":       "浅色步骤",
    "quote":         "深色引用",
    "end":           "深色结束",
}

def _check(): return _HAS_PPTX

def _rgb(h):
    from pptx.dml.color import RGBColor
    return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

def _In(n):
    from pptx.util import Inches; return Inches(n)

def _pt(n):
    from pptx.util import Pt; return Pt(n)

def _add_rect(slide, x, y, w, h, fill):
    from pptx.util import Inches
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill)
    s.line.fill.background()
    return s

def _add_text(slide, text, x, y, w, h, size=18, bold=False, color="FFFFFF", font="Calibri", align="left", italic=False):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    if not text: text = ""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.paragraphs[0].alignment = {"left":PP_ALIGN.LEFT,"center":PP_ALIGN.CENTER,"right":PP_ALIGN.RIGHT}[align]
    r = tf.paragraphs[0].add_run()
    r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = _rgb(color); r.font.name = font
    r.font.italic = italic
    return tb

def _add_bullets(slide, items, x, y, w, h, size=16, color="333333", spacing=8):
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    if isinstance(items, str): items = [items]
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        text = str(item.get("text", item.get("label", str(item)))) if isinstance(item, dict) else str(item)
        r = p.add_run()
        r.text = "  " + text
        r.font.size = Pt(size); r.font.color.rgb = _rgb(color); r.font.name = "Calibri"
    return tb

# ── Design constants ──
C = {
    "dark_bg": "0F172A", "light_bg": "F8FAFC",
    "blue": "3B82F6", "indigo": "6366F1", "gold": "F59E0B",
    "white": "FFFFFF", "near_black": "111827", "gray": "6B7280",
    "light_gray": "9CA3AF", "green": "10B981", "red": "EF4444",
    "card_white": "FFFFFF", "card_dark": "1E293B",
}

def _title_slide(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(sl, 0, 0, 13.33, 7.5, C["dark_bg"])
    _add_rect(sl, 0.8, 2.85, 1.8, 0.04, C["gold"])
    _add_text(sl, data.get("title",""), 0.8, 1.3, 10.5, 1.5, 44, True, C["white"], "Georgia")
    _add_text(sl, data.get("subtitle",""), 0.8, 3.2, 10.5, 1.5, 20, False, C["light_gray"], "Calibri")

def _title_center_slide(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(sl, 0, 0, 13.33, 7.5, C["dark_bg"])
    _add_text(sl, data.get("title",""), 1.0, 2.0, 11.33, 1.2, 40, True, C["gold"], "Georgia", "center")
    _add_rect(sl, 5.9, 3.3, 1.5, 0.03, C["blue"])
    _add_text(sl, data.get("subtitle",""), 1.5, 3.6, 10.33, 1.5, 18, False, C["light_gray"], "Calibri", "center")

def _content_light_slide(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(sl, 0, 0, 13.33, 7.5, C["light_bg"])
    _add_rect(sl, 0, 0, 13.33, 0.05, C["blue"])
    _add_text(sl, data.get("title",""), 0.8, 0.35, 10, 0.8, 30, True, C["near_black"], "Georgia")
    _add_rect(sl, 0.8, 1.2, 1.6, 0.03, C["gold"])
    content = data.get("content", data.get("text",""))
    if isinstance(content, list): _add_bullets(sl, content, 0.8, 1.6, 10.5, 4.5, 18, C["gray"])
    else: _add_text(sl, content, 0.8, 1.6, 10.5, 4.5, 18, False, C["near_black"])

def _content_dark_slide(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(sl, 0, 0, 13.33, 7.5, C["dark_bg"])
    _add_text(sl, data.get("title",""), 0.8, 0.35, 10, 0.8, 30, True, C["white"], "Georgia")
    _add_rect(sl, 0.8, 1.2, 1.6, 0.03, C["gold"])
    content = data.get("content", data.get("text",""))
    if isinstance(content, list): _add_bullets(sl, content, 0.8, 1.6, 10.5, 4.5, 18, C["light_gray"])
    else: _add_text(sl, content, 0.8, 1.6, 10.5, 4.5, 18, False, C["light_gray"])

def _two_column_slide(prs, data, dark=False):
    bg = C["dark_bg"] if dark else C["light_bg"]
    txt = C["white"] if dark else C["near_black"]
    sub = C["light_gray"] if dark else C["gray"]
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(sl, 0, 0, 13.33, 7.5, bg)
    if not dark: _add_rect(sl, 0, 0, 13.33, 0.05, C["blue"])
    _add_text(sl, data.get("title",""), 0.8, 0.35, 10, 0.8, 30, True, txt, "Georgia")
    _add_rect(sl, 0.8, 1.2, 1.6, 0.03, C["gold"])
    # Left column
    _add_rect(sl, 0.8, 1.5, 5.5, 0.04, C["blue"])
    c1 = data.get("content", data.get("text",""))
    if isinstance(c1, list): _add_bullets(sl, c1, 0.8, 1.8, 5.3, 4.0, 16, sub)
    else: _add_text(sl, c1, 0.8, 1.8, 5.3, 4.0, 16, False, sub)
    # Right column
    _add_rect(sl, 7.0, 1.5, 5.5, 0.04, C["green"])
    c2 = data.get("content2","")
    if isinstance(c2, list): _add_bullets(sl, c2, 7.0, 1.8, 5.3, 4.0, 16, sub)
    else: _add_text(sl, c2, 7.0, 1.8, 5.3, 4.0, 16, False, sub)

def _features_slide(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(sl, 0, 0, 13.33, 7.5, C["dark_bg"])
    _add_text(sl, data.get("title",""), 0.8, 0.35, 10, 0.8, 30, True, C["white"], "Georgia")
    _add_rect(sl, 0.8, 1.2, 1.6, 0.03, C["gold"])
    items = data.get("content", data.get("text", []))
    if isinstance(items, str): items = [items]
    colors = [C["blue"], C["green"], C["gold"], C["indigo"], C["red"], "EC4899"]
    for i, it in enumerate(items):
        col = i % 3; row = i // 3
        cx = 0.8 + col * 4.0; cy = 1.6 + row * 2.3
        label = it.get("label", str(it)) if isinstance(it, dict) else str(it)
        desc = it.get("desc", "") if isinstance(it, dict) else ""
        # Colored top rule instead of card background
        _add_rect(sl, cx, cy, 2.4, 0.04, colors[i % len(colors)])
        _add_text(sl, label, cx, cy + 0.2, 3.4, 0.55, 20, True, C["white"], "Calibri")
        if desc: _add_text(sl, desc, cx, cy + 0.8, 3.4, 1.0, 14, False, C["light_gray"], "Calibri")

def _process_slide(prs, data, dark=False):
    bg = C["dark_bg"] if dark else C["light_bg"]
    txt = C["white"] if dark else C["near_black"]
    sub = C["light_gray"] if dark else C["gray"]
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(sl, 0, 0, 13.33, 7.5, bg)
    if not dark: _add_rect(sl, 0, 0, 13.33, 0.05, C["blue"])
    _add_text(sl, data.get("title",""), 0.8, 0.35, 10, 0.8, 30, True, txt, "Georgia")
    _add_rect(sl, 0.8, 1.2, 1.6, 0.03, C["gold"])
    items = data.get("content", data.get("text", []))
    if isinstance(items, str): items = [items]
    colors = [C["blue"], C["indigo"], C["green"], "EC4899", C["gold"]]
    for i, it in enumerate(items):
        label = it.get("label", str(it)) if isinstance(it, dict) else str(it)
        desc = it.get("desc", "") if isinstance(it, dict) else ""
        cy = 1.6 + i * 1.15
        _add_rect(sl, 0.8, cy + 0.05, 0.45, 0.45, colors[i % len(colors)])
        _add_text(sl, str(i+1), 0.8, cy + 0.08, 0.45, 0.4, 22, True, C["white"], "Calibri", "center")
        _add_text(sl, label, 1.55, cy + 0.02, 10.5, 0.45, 18, True, txt, "Calibri")
        if desc: _add_text(sl, desc, 1.55, cy + 0.45, 10.5, 0.5, 14, False, sub, "Calibri")

def _quote_slide(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(sl, 0, 0, 13.33, 7.5, C["dark_bg"])
    _add_rect(sl, 1.5, 2.6, 10.33, 0.03, C["gold"])
    _add_text(sl, data.get("title",""), 1.5, 1.5, 10.33, 1.5, 30, False, C["gold"], "Georgia", italic=True)
    _add_text(sl, data.get("subtitle",""), 1.5, 3.0, 10.33, 0.8, 16, False, C["light_gray"], "Calibri", "right")

def _end_slide(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(sl, 0, 0, 13.33, 7.5, C["dark_bg"])
    _add_rect(sl, 0, 7.42, 13.33, 0.08, C["gold"])
    _add_text(sl, data.get("title","谢谢"), 0.8, 2.0, 10.5, 1.5, 56, True, C["white"], "Georgia")
    _add_rect(sl, 0.8, 3.6, 2.0, 0.04, C["gold"])
    _add_text(sl, data.get("subtitle",""), 0.8, 3.9, 10.5, 1.0, 18, False, C["light_gray"], "Calibri")

_BUILDERS = {
    "title": _title_slide, "title_dark": _title_slide,
    "title_center": _title_center_slide,
    "content_light": _content_light_slide, "content": _content_light_slide,
    "content_dark": _content_dark_slide,
    "two_column": lambda p,d: _two_column_slide(p,d,False),
    "two_column_light": lambda p,d: _two_column_slide(p,d,False),
    "two_column_dark": lambda p,d: _two_column_slide(p,d,True),
    "features_dark": _features_slide, "features": _features_slide,
    "process_light": lambda p,d: _process_slide(p,d,False), "process": lambda p,d: _process_slide(p,d,False),
    "process_dark": lambda p,d: _process_slide(p,d,True),
    "quote_dark": _quote_slide, "quote": _quote_slide,
    "end": _end_slide, "end_dark": _end_slide,
}

# ── 纯 Python 回退引擎 ──
def _build_pure(data, layout, pb):
    C2 = {"dark_bg":"0F172A","light_bg":"F8FAFC","blue":"3B82F6","gold":"F59E0B",
          "white":"FFFFFF","near_black":"111827","gray":"6B7280","light_gray":"9CA3AF","card_dark":"1E293B"}
    title = data.get("title","")
    subtitle = data.get("subtitle","")
    content = data.get("content", data.get("text",""))
    if isinstance(content, dict) and "items" in content: content = content["items"]
    if not isinstance(content, list): content = [content] if content else []

    if layout in ("title","title_dark"):
        sl = pb.add_slide(C2["dark_bg"])
        pb.add_text(sl, title, 0.8, 2.6, 10.5, 0.8, 36, True, C2["white"], "Georgia")
        pb.add_rect(sl, 0.8, 3.5, 1.8, 0.03, C2["gold"])
        pb.add_text(sl, subtitle, 0.8, 3.8, 10.5, 1.0, 20, False, C2["light_gray"])
    elif layout == "title_center":
        sl = pb.add_slide(C2["dark_bg"])
        pb.add_text(sl, title, 1.0, 2.5, 11.33, 0.8, 36, True, C2["gold"], "Georgia", "center")
        pb.add_rect(sl, 5.8, 3.4, 1.5, 0.03, C2["blue"])
        pb.add_text(sl, subtitle, 1.5, 3.7, 10.33, 1.0, 18, False, C2["light_gray"], "Calibri", "center")
    elif layout in ("content_light","content"):
        sl = pb.add_slide(C2["light_bg"])
        pb.add_rect(sl, 0, 0, 13.33, 0.04, C2["blue"])
        pb.add_text(sl, title, 0.8, 0.35, 10, 0.7, 28, True, C2["near_black"], "Georgia")
        pb.add_rect(sl, 0.8, 1.1, 1.6, 0.03, C2["gold"])
        pb.add_bullets(sl, content, 0.8, 1.5, 10.5, 4.5, 18, C2["gray"])
    elif layout == "content_dark":
        sl = pb.add_slide(C2["dark_bg"])
        pb.add_text(sl, title, 0.8, 0.35, 10, 0.7, 28, True, C2["white"], "Georgia")
        pb.add_rect(sl, 0.8, 1.1, 1.6, 0.03, C2["gold"])
        pb.add_bullets(sl, content, 0.8, 1.5, 10.5, 4.5, 18, C2["light_gray"])
    elif layout in ("two_column","two_column_light"):
        sl = pb.add_slide(C2["light_bg"])
        pb.add_rect(sl, 0, 0, 13.33, 0.04, C2["blue"])
        pb.add_text(sl, title, 0.8, 0.35, 10, 0.7, 28, True, C2["near_black"], "Georgia")
        pb.add_rect(sl, 0.8, 1.1, 1.6, 0.03, C2["gold"])
        pb.add_rect(sl, 0.8, 1.4, 5.5, 0.04, C2["blue"])
        c2 = data.get("content2","")
        c2 = [c2] if isinstance(c2,str) and c2 else (c2 if isinstance(c2,list) else [])
        pb.add_bullets(sl, content, 0.8, 1.7, 5.3, 4.0, 16, C2["near_black"])
        pb.add_rect(sl, 7.0, 1.4, 5.5, 0.04, C2["gold"])
        pb.add_bullets(sl, c2, 7.0, 1.7, 5.3, 4.0, 16, C2["near_black"])
    elif layout == "two_column_dark":
        sl = pb.add_slide(C2["dark_bg"])
        pb.add_text(sl, title, 0.8, 0.35, 10, 0.7, 28, True, C2["white"], "Georgia")
        pb.add_rect(sl, 0.8, 1.1, 1.6, 0.03, C2["gold"])
        pb.add_rect(sl, 0.8, 1.4, 5.5, 0.04, C2["blue"])
        c2 = data.get("content2","")
        c2 = [c2] if isinstance(c2,str) and c2 else (c2 if isinstance(c2,list) else [])
        pb.add_bullets(sl, content, 0.8, 1.7, 5.3, 4.0, 16, C2["light_gray"])
        pb.add_rect(sl, 7.0, 1.4, 5.5, 0.04, C2["gold"])
        pb.add_bullets(sl, c2, 7.0, 1.7, 5.3, 4.0, 16, C2["light_gray"])
    elif layout in ("features_dark","features"):
        sl = pb.add_slide(C2["dark_bg"])
        pb.add_text(sl, title, 0.8, 0.35, 10, 0.7, 28, True, C2["white"], "Georgia")
        pb.add_rect(sl, 0.8, 1.1, 1.6, 0.03, C2["gold"])
        pb.add_bullets(sl, content, 0.8, 1.5, 10.5, 4.5, 18, C2["light_gray"])
    elif layout in ("process_light","process"):
        sl = pb.add_slide(C2["light_bg"])
        pb.add_rect(sl, 0, 0, 13.33, 0.04, C2["blue"])
        pb.add_text(sl, title, 0.8, 0.35, 10, 0.7, 28, True, C2["near_black"], "Georgia")
        pb.add_rect(sl, 0.8, 1.1, 1.6, 0.03, C2["gold"])
        pb.add_bullets(sl, content, 0.8, 1.5, 10.5, 4.5, 18, C2["gray"])
    elif layout in ("quote_dark","quote"):
        sl = pb.add_slide(C2["dark_bg"])
        pb.add_text(sl, title, 1.5, 2.0, 10.33, 1.2, 28, False, C2["gold"], "Georgia", italic=True)
        pb.add_rect(sl, 1.5, 3.3, 10.33, 0.03, C2["gold"])
        pb.add_text(sl, subtitle, 1.5, 3.6, 10.33, 0.8, 16, False, C2["light_gray"], "Calibri", "right")
    elif layout in ("end","end_dark"):
        sl = pb.add_slide(C2["dark_bg"])
        pb.add_text(sl, title or "Thanks", 0.8, 2.5, 10.5, 1.0, 48, True, C2["white"], "Georgia")
        pb.add_rect(sl, 0.8, 3.6, 2.0, 0.04, C2["gold"])
        pb.add_text(sl, subtitle, 0.8, 3.9, 10.5, 1.0, 18, False, C2["light_gray"])

@dataclass
class CreateModernPPTXTool:
    @property
    def definition(self) -> ToolDefinition:
        layouts = "、".join(f"{k}({v})" for k,v in LAYOUTS.items())
        engines = "pure(纯Python标准库,零依赖,推荐) 或 pptx(python-pptx库,更多样式)" if _HAS_PPTX else "pure(纯Python标准库,零依赖)"
        return ToolDefinition(
            name="create_modern_pptx",
            description=f"创建精美PPT（极简现代风格）。可用引擎: {engines}。布局: {layouts}。示例: [{{\"type\":\"title\",\"title\":\"产品发布\",\"subtitle\":\"2026年\"}},{{\"type\":\"two_column\",\"title\":\"对比\",\"content\":[\"优势A\"],\"content2\":[\"劣势B\"]}},{{\"type\":\"end\",\"title\":\"谢谢\"}}]",
            parameters={"type":"object","properties":{
                "slides":{"type":"string","description":"JSON幻灯片数组"},
                "engine":{"type":"string",f"description":"生成引擎: {engines}","default":"pure"},
                "output":{"type":"string","description":"输出路径","default":""}
            }}
        )

    async def execute(self, params: Dict[str, Any]) -> ToolResult:
        try:
            raw = params.get("slides","[]")
            if isinstance(raw, str):
                s = raw.strip()
                if s.startswith("```"): s = s.split("\n",1)[1].rsplit("\n```",1)[0].strip()
                slides = json.loads(s)
            else:
                slides = raw

            engine = params.get("engine","pure")

            output = params.get("output","").strip()
            if not output:
                ws = os.path.join(os.path.dirname(os.path.dirname(__file__)),"workspace")
                os.makedirs(ws, exist_ok=True)
                output = os.path.join(ws, f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
            output = os.path.abspath(output)

            if engine=="pptx" and _HAS_PPTX:
                # ── python-pptx 引擎（用户显式指定） ──
                prs = Presentation()
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                for sd in slides:
                    lt = sd.get("type","content_light")
                    _BUILDERS.get(lt, _content_light_slide)(prs, sd)
                prs.save(output)
            else:
                # ── 纯 Python 引擎（默认/推荐） ──
                from .pttx_builder import PttxBuilder
                pb = PttxBuilder()
                for sd in slides:
                    lt = sd.get("type","content_light")
                    _build_pure(sd, lt, pb)
                pb.save(output)

            return ToolResult(success=True,
                content=f"✅ PPT生成成功！\n📄 {output}\n📑 {len(slides)}页\n引擎: {'python-pptx' if (engine=='pptx' and _HAS_PPTX) else '纯Python'}\n布局: {[s.get('type','?') for s in slides]}")

        except json.JSONDecodeError as e:
            return ToolResult(success=False,content="",error=f"JSON错误: {e}")
        except Exception as e:
            return ToolResult(success=False,content="",error=str(e))


# ── 旧工具 ──
THEME_COLORS = {"blue":0x1A365D,"green":0x22543D,"red":0x742A2A,"purple":0x44337A,"orange":0x7B341E,"cyan":0x0E6B6B,"gray":0x2D3748}

def _theme(slide, tn):
    from pptx.dml.color import RGBColor
    c = RGBColor((THEME_COLORS.get(tn,0x1A365D)>>16)&0xFF,(THEME_COLORS.get(tn,0x1A365D)>>8)&0xFF,THEME_COLORS.get(tn,0x1A365D)&0xFF)
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c
    if slide.shapes.title:
        for p in slide.shapes.title.text_frame.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

def _ph(s,i,t):
    try:
        if i < len(s.placeholders): s.placeholders[i].text = t; return True
    except: return False

@dataclass
class CreatePPTXTool:
    @property
    def definition(self): return ToolDefinition(name="create_pptx",description="基础PPT（内置布局）",parameters={"type":"object","properties":{"title":{"type":"string","default":""},"subtitle":{"type":"string","default":""},"theme":{"type":"string","default":"blue"},"output":{"type":"string","default":""}}})
    async def execute(self,params):
        if not _check(): return ToolResult(success=False,content="",error="pip install python-pptx")
        try:
            from pptx import Presentation; p = Presentation()
            s = p.slides.add_slide(p.slide_layouts[0])
            if s.shapes.title: s.shapes.title.text = str(params.get("title","")).strip()
            _ph(s,1,str(params.get("subtitle","")).strip()); _theme(s,params.get("theme","blue"))
            o = params.get("output","").strip() or f"ppt_{datetime.now().strftime('%H%M%S')}.pptx"
            p.save(o); return ToolResult(success=True,content=f"✅ {os.path.abspath(o)}")
        except Exception as e: return ToolResult(success=False,content="",error=str(e))

@dataclass
class CreateSmartPPTXTool:
    @property
    def definition(self): return ToolDefinition(name="create_smart_ppt",description="智能PPT（基础版）",parameters={"type":"object","properties":{"title":{"type":"string","default":""},"sections":{"type":"string","default":""},"theme":{"type":"string","default":"blue"},"output":{"type":"string","default":""}}})
    async def execute(self,params):
        if not _check(): return ToolResult(success=False,content="",error="pip install python-pptx")
        try:
            ss = _parse(params.get("sections",""))
            if not ss: return ToolResult(success=False,content="",error="内容不能为空")
            from pptx import Presentation; p = Presentation()
            s = p.slides.add_slide(p.slide_layouts[0])
            if s.shapes.title: s.shapes.title.text = str(params.get("title","")).strip()
            _ph(s,1,f"共 {len(ss)} 章节"); _theme(s,params.get("theme","blue"))
            for sec in ss:
                sl = p.slides.add_slide(p.slide_layouts[1])
                if sl.shapes.title: sl.shapes.title.text = str(sec.get("title",sec.get("label","")))
                _ph(sl,1,str(sec.get("content",sec.get("desc",sec.get("details",""))))); _theme(sl,params.get("theme","blue"))
            o = params.get("output","").strip() or f"smart_{datetime.now().strftime('%H%M%S')}.pptx"
            p.save(o); return ToolResult(success=True,content=f"✅ {os.path.abspath(o)}")
        except Exception as e: return ToolResult(success=False,content="",error=str(e))

def _parse(t):
    t = t.strip()
    if t.startswith("[") or t.startswith("{"):
        try: d = json.loads(t); return d if isinstance(d,list) else [d]
        except: pass
    return [{"title":l.split("：" if "：" in l else ":",1)[0].strip(),"content":l.split("：" if "：" in l else ":",1)[1].strip() if "：" in l or ":" in l else ""} for l in [l.strip() for l in t.split('\n') if l.strip()]]

class PPTSkill:
    @property
    def metadata(self): return SkillMetadata(name="PPT制作",description="极简现代风格PPT：10种布局、纯Python(零Node.js)。create_modern_pptx(推荐)/create_pptx/create_smart_ppt",author="PyClaw",version="4.0.0",tags=["ppt","pptx","python-pptx"])
    def get_tools(self): return [CreateModernPPTXTool(),CreatePPTXTool(),CreateSmartPPTXTool()]
    async def initialize(self)->bool:
        print("✅ PPT Skill v4.0 (极简现代) 已加载" if _check() else "⚠️ pip install python-pptx")
        return True
    async def cleanup(self): print("📊 PPT Skill 已卸载")
